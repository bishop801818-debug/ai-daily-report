"""
法恩莱特撬装工厂战略规划 PPT 生成脚本
专业版 - 包含美化设计、统一格式、专业配色
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# 定义专业配色方案
class Colors:
    PRIMARY = RGBColor(16, 78, 139)      # 深蓝色 - 主色
    SECONDARY = RGBColor(230, 81, 0)     # 橙红色 - 强调色
    ACCENT = RGBColor(0, 150, 136)       # 青绿色 - 辅助色
    SUCCESS = RGBColor(76, 175, 80)      # 绿色 - 成功
    WARNING = RGBColor(255, 152, 0)      # 橙色 - 警告
    DANGER = RGBColor(244, 67, 54)       # 红色 - 危险
    LIGHT_BG = RGBColor(245, 247, 250)   # 浅灰背景
    DARK_TEXT = RGBColor(33, 33, 33)     # 深灰文字
    MEDIUM_TEXT = RGBColor(117, 117, 117) # 中灰文字
    WHITE = RGBColor(255, 255, 255)      # 白色

# 定义字体
class Fonts:
    TITLE = '微软雅黑'
    CONTENT = '微软雅黑'
    CODE = 'Consolas'

def set_slide_background(slide, prs, gradient=False):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    
    if gradient:
        # 渐变背景
        fill.gradient()
        fill.gradient_angle = 90
        fill.gradient_stops[0].color.rgb = Colors.PRIMARY
        fill.gradient_stops[1].color.rgb = RGBColor(5, 50, 100)
    else:
        # 纯色背景
        fill.solid()
        fill.fore_color.rgb = Colors.WHITE

def add_title_slide(prs, title, subtitle):
    """添加专业标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 设置背景
    set_slide_background(slide, prs, gradient=True)
    
    # 添加标题背景条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.5),
        Inches(13.333), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.WHITE
    shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = Colors.PRIMARY
    title_p.font.name = Fonts.TITLE
    title_p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = Colors.DARK_TEXT
    subtitle_p.font.name = Fonts.CONTENT
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 添加装饰线条
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(4.3),
        Inches(3.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Colors.SECONDARY
    line.line.fill.background()
    
    return slide

def add_section_slide(prs, section_title):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景
    set_slide_background(slide, prs, gradient=True)
    
    # 添加半透明遮罩
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = Colors.PRIMARY
    overlay.fill.transparency = 0.1
    overlay.line.fill.background()
    
    # 添加章节标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = section_title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = Colors.WHITE
    title_p.font.name = Fonts.TITLE
    title_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, highlight_first=False):
    """添加专业内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景
    set_slide_background(slide, prs)
    
    # 添加顶部色条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.6)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = Colors.PRIMARY
    header_bar.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.4))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = Colors.WHITE
    title_p.font.name = Fonts.TITLE
    
    # 添加页码装饰
    page_num = slide.shapes.add_textbox(Inches(12), Inches(0.2), Inches(1), Inches(0.3))
    page_num.text_frame.text = '●'
    page_num.text_frame.paragraphs[0].font.color.rgb = Colors.SECONDARY
    page_num.text_frame.paragraphs[0].font.size = Pt(20)
    
    # 添加内容
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(12), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            # 处理缩进
            if item.startswith('  '):
                p.text = item.strip()
                p.font.size = Pt(18)
                p.level = 1
                p.space_before = Pt(6)
            elif item.startswith('●') or item.startswith('◆') or item.startswith('✓'):
                p.text = item
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = Colors.PRIMARY
                p.space_before = Pt(12)
                p.space_after = Pt(6)
            else:
                p.text = item
                p.font.size = Pt(20)
                p.font.color.rgb = Colors.DARK_TEXT
                p.space_before = Pt(8)
                p.space_after = Pt(4)
            
            p.font.name = Fonts.CONTENT
    
    return slide

def add_table_slide(prs, title, table_data, highlight_header=True):
    """添加专业表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景
    set_slide_background(slide, prs)
    
    # 添加顶部色条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.6)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = Colors.PRIMARY
    header_bar.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.4))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = Colors.WHITE
    title_p.font.name = Fonts.TITLE
    
    # 添加表格
    if table_data:
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows > 0 and cols > 0:
            # 计算表格尺寸
            table_width = Inches(12)
            table_height = Inches(0.7 * rows)
            table_x = Inches(0.65)
            table_y = Inches(0.9)
            
            table_shape = slide.shapes.add_table(
                rows, cols,
                table_x, table_y,
                table_width, table_height
            )
            
            table = table_shape.table
            
            # 设置列宽
            col_width = int(table_width / cols)
            for i in range(cols):
                table.columns[i].width = col_width
            
            # 填充数据
            for i, row in enumerate(table_data):
                for j, cell_text in enumerate(row):
                    cell = table.cell(i, j)
                    cell.text = str(cell_text)
                    cell.text_frame.paragraphs[0].font.size = Pt(14)
                    cell.text_frame.paragraphs[0].font.name = Fonts.CONTENT
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # 设置单元格边距
                    cell.margin_top = Inches(0.1)
                    cell.margin_bottom = Inches(0.1)
                    cell.margin_left = Inches(0.05)
                    cell.margin_right = Inches(0.05)
                    
                    # 设置表头样式
                    if i == 0 and highlight_header:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = Colors.PRIMARY
                        cell.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
                        cell.text_frame.paragraphs[0].font.bold = True
                    else:
                        # 设置交替行颜色
                        if i % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = Colors.LIGHT_BG
                        cell.text_frame.paragraphs[0].font.color.rgb = Colors.DARK_TEXT
    
    return slide

def add_comparison_slide(prs, title, comparison_data):
    """添加专业对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景
    set_slide_background(slide, prs)
    
    # 添加顶部色条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.6)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = Colors.PRIMARY
    header_bar.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.4))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = Colors.WHITE
    title_p.font.name = Fonts.TITLE
    
    # 添加对比内容
    if comparison_data:
        y_position = 0.9
        for i, item in enumerate(comparison_data):
            # 添加项目符号
            if item.startswith('●') or item.startswith('◆') or item.startswith('✓'):
                textbox = slide.shapes.add_textbox(Inches(0.7), Inches(y_position), Inches(12), Inches(0.5))
                textbox.text_frame.text = item
                textbox.text_frame.paragraphs[0].font.size = Pt(20)
                textbox.text_frame.paragraphs[0].font.name = Fonts.CONTENT
                
                # 设置颜色
                if '优势' in item or '成功' in item or '✅' in item:
                    textbox.text_frame.paragraphs[0].font.color.rgb = Colors.SUCCESS
                elif '风险' in item or '警告' in item:
                    textbox.text_frame.paragraphs[0].font.color.rgb = Colors.WARNING
                else:
                    textbox.text_frame.paragraphs[0].font.color.rgb = Colors.DARK_TEXT
                
                y_position += 0.5
            else:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(11.7), Inches(0.4))
                textbox.text_frame.text = item
                textbox.text_frame.paragraphs[0].font.size = Pt(18)
                textbox.text_frame.paragraphs[0].font.name = Fonts.CONTENT
                textbox.text_frame.paragraphs[0].font.color.rgb = Colors.MEDIUM_TEXT
                y_position += 0.45
    
    return slide

def add_chart_placeholder_slide(prs, title, chart_description, data_points=None):
    """添加图表占位页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景
    set_slide_background(slide, prs)
    
    # 添加顶部色条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.6)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = Colors.PRIMARY
    header_bar.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.4))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = Colors.WHITE
    title_p.font.name = Fonts.TITLE
    
    # 添加图表占位框
    chart_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(1.5),
        Inches(9.333), Inches(4.5)
    )
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = Colors.LIGHT_BG
    chart_box.line.width = Pt(2)
    chart_box.line.color.rgb = Colors.PRIMARY
    chart_box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    
    # 添加图表说明
    desc_box = slide.shapes.add_textbox(Inches(2.2), Inches(3.5), Inches(9), Inches(2))
    desc_frame = desc_box.text_frame
    desc_p = desc_frame.paragraphs[0]
    desc_p.text = chart_description
    desc_p.font.size = Pt(16)
    desc_p.font.color.rgb = Colors.MEDIUM_TEXT
    desc_p.font.name = Fonts.CONTENT
    desc_p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_ppt():
    """生成专业 PPT"""
    print("开始生成专业 PPT...")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置页面大小为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 封面页 ==========
    print("添加封面页...")
    add_title_slide(
        prs,
        "法恩莱特撬装工厂战略规划",
        "宁德时代山东、河南基地就近配套项目\n2026 年 4 月"
    )
    
    # ========== 目录页 ==========
    print("添加目录页...")
    add_content_slide(prs, "目录", [
        "● 战略机遇验证",
        "  宁德时代贴厂模式调研",
        "  山东、河南基地布局分析",
        "● 战略定位与目标",
        "  核心价值主张",
        "  短中长期目标",
        "● 技术方案",
        "  撬装工厂设计",
        "  JIT 供应流程",
        "● 投资预算与回报",
        "  投资明细",
        "  收益分析",
        "● 实施路径",
        "  四阶段推进计划",
        "● 核心竞争优势",
        "  五大优势分析",
        "● 行动计划",
        "  下一步工作安排"
    ])
    
    # ========== 第一部分：战略机遇验证 ==========
    print("添加第一部分...")
    add_section_slide(prs, "第一部分\n战略机遇验证")
    
    # 执行摘要
    add_content_slide(prs, "执行摘要 - 核心结论", [
        "✅ 宁德时代贴厂模式已验证",
        "  重庆'厂中厂'模式 2025 年 6 月投产",
        "  电池下线 1 小时内装车",
        "  供应链响应从数周压缩到分钟级",
        "",
        "◆ 山东、河南基地真实存在且规模巨大",
        "  山东济宁：2026 年产能超 100GWh",
        "  河南洛阳：2026 年产能超 100GWh",
        "",
        "✓ 电解液需求巨大",
        "  两大基地年需求 20-24 万吨",
        "  市场空间 60-96 亿元/年",
        "",
        "◆ 撬装工厂模式完全可行",
        "  投资回收期 2-3 年",
        "  建议立即启动"
    ])
    
    # 宁德时代贴厂模式
    add_content_slide(prs, "宁德时代贴厂模式 - 重庆'厂中厂'案例", [
        "◆ 项目概况",
        "  时间：2025 年 6 月 30 日正式投产",
        "  地点：重庆赛力斯超级工厂内部",
        "  产能：两条 CTP2.0 高端电池包产线",
        "",
        "✓ 核心数据",
        "  电池下线后不到 1 小时就能装车",
        "  从下达指令到产线调整只需不到 20 分钟",
        "  每天数千个零部件当天周转完毕",
        "  物流响应周期从数天/数周压缩到分钟级",
        "",
        "◆ 模式特点",
        "  物理空间融合：供应商车间直接嵌入整车厂内部",
        "  生产系统协同：共享同一套生产节奏和数据系统",
        "  零距离供应：取消中间物流环节",
        "  零时差响应：实时响应生产计划变化"
    ])
    
    # 山东基地布局
    add_table_slide(prs, "宁德时代山东基地布局", [
        ["项目", "详情"],
        ["位置", "山东省济宁市兖州区"],
        ["投资方", "宁德时代全资子公司 - 山东时代"],
        ["投产时间", "2025 年 5 月 17 日一期投产"],
        ["战略定位", "北方首个基地，北方产能规划最大基地"],
        ["一期产能", "60GWh，占地 800 余亩，投资 140 亿元"],
        ["二期产能", "30GWh，2025-2026 年"],
        ["三期产能", "40GWh+，2026 年"],
        ["总产能 (2026)", "超 100GWh，总占地约 2000 亩"],
        ["总投资", "超 210 亿元"],
        ["建设标准", "灯塔工厂 + 零碳工厂"]
    ])
    
    # 河南基地布局
    add_table_slide(prs, "宁德时代河南基地布局", [
        ["项目", "详情"],
        ["位置", "河南省洛阳市伊滨区"],
        ["开工时间", "2022 年 9 月一期开工"],
        ["一期投产", "2024 年 7 月电池包工厂，2025 年 4 月全线投产"],
        ["二期投产", "2025 年 9 月/10 月投产"],
        ["一期产能", "30GWh，占地 850 亩"],
        ["二期产能", "30GWh"],
        ["三四期", "施工中，2026 年投产"],
        ["总产能", "超 100GWh，总占地 3400 亩"],
        ["总投资", "超 300 亿元"],
        ["产业链配套", "26 个上下游项目签约落地"],
        ["年产值", "全部建成后超千亿元"]
    ])
    
    # 电解液需求测算
    add_table_slide(prs, "电解液需求测算", [
        ["基地", "2026 年产能", "电解液年需求", "市场空间"],
        ["山东济宁", "100GWh", "10-12 万吨", "30-48 亿元/年"],
        ["河南洛阳", "100GWh", "10-12 万吨", "30-48 亿元/年"],
        ["合计", "200GWh", "20-24 万吨", "60-96 亿元/年"],
        [""],
        ["计算公式", "电解液需求（吨）= 电池产能（GWh）× 1000-1200 吨/GWh", ""]
    ])
    
    # ========== 第二部分：战略定位与目标 ==========
    print("添加第二部分...")
    add_section_slide(prs, "第二部分\n战略定位与目标")
    
    # 战略定位
    add_content_slide(prs, "战略定位", [
        "◆ 总体定位",
        "  成为宁德时代山东、河南基地的核心电解液战略供应商，",
        "  采用撬装工厂模式实现就近配套、JIT 供应，",
        "  打造零距离、零时差、零碳排的供应链标杆",
        "",
        "✓ 核心价值主张 - 四个'零'",
        "  零距离：撬装工厂部署在宁德时代基地周边 50 公里范围内",
        "  零时差：管道输送 + 自动化产线，响应时间<2 小时",
        "  零碳排：配套光伏 + 储能，使用绿电生产",
        "  零库存：按订单生产，JIT 供应",
        "",
        "◆ 三大优势",
        "  轻资产：模块化设计，可快速复制、灵活迁移",
        "  快响应：建设周期 3-6 个月，远快于传统工厂 2-3 年",
        "  低成本：投资成本降低 50%，运营成本降低 30%"
    ])
    
    # 战略目标体系
    add_content_slide(prs, "战略目标体系", [
        "◆ 短期目标（2026 年）",
        "  完成 1 座撬装工厂示范项目建设",
        "  实现 5GWh 配套产能（约 5000 吨/年）",
        "  通过宁德时代供应商认证",
        "  示范工厂投资 5000-8000 万元",
        "",
        "✓ 中期目标（2027-2028 年）",
        "  在山东济宁、河南洛阳各部署 1 座撬装工厂",
        "  总产能达到 20GWh 配套（约 2 万吨/年）",
        "  进入宁德时代核心供应商体系",
        "  年销售收入达到 15-20 亿元",
        "",
        "◆ 长期目标（2029-2030 年）",
        "  部署 4-6 座撬装工厂，覆盖宁德时代全国基地",
        "  总产能达到 50GWh 配套（约 5 万吨/年）",
        "  成为宁德时代前三大电解液供应商",
        "  年销售收入达到 40-50 亿元"
    ])
    
    # ========== 第三部分：技术方案 ==========
    print("添加第三部分...")
    add_section_slide(prs, "第三部分\n技术方案")
    
    # 撬装工厂设计参数
    add_table_slide(prs, "撬装工厂设计参数对比", [
        ["参数项", "指标", "传统工厂对比"],
        ["占地面积", "800-1000㎡", "传统工厂 1/10"],
        ["建筑面积", "1000-1500㎡", "传统工厂 1/8"],
        ["建设周期", "3-6 个月", "传统工厂 1/3"],
        ["投资规模", "5000-8000 万元", "传统工厂 1/2"],
        ["产能规模", "5-10GWh 配套", "可灵活扩展"],
        ["人员配置", "10-15 人", "传统工厂 1/5"],
        ["自动化程度", "L4 级（无人化）", "行业领先"]
    ])
    
    # 核心功能模块
    add_content_slide(prs, "撬装工厂核心功能模块", [
        "◆ 模块一：原料储存模块",
        "  NMP 储罐、锂盐储罐、添加剂储罐",
        "  氮气保护系统、温湿度控制",
        "",
        "✓ 模块二：配料混合模块",
        "  高精度计量系统（精度±0.1%）",
        "  自动配料系统、真空搅拌系统",
        "",
        "◆ 模块三：过滤灌装模块",
        "  多级过滤系统（精度 1μm）",
        "  自动灌装系统、密封检测",
        "",
        "✓ 模块四：成品仓储模块",
        "  恒温恒湿仓库、自动立体仓库",
        "  AGV 搬运系统、WMS 管理系统",
        "",
        "◆ 模块五：管道输送模块（可选）",
        "  不锈钢管道系统、计量泵组",
        "  直接输送到宁德时代产线"
    ])
    
    # JIT 供应流程
    add_content_slide(prs, "JIT 供应流程与响应时间", [
        "◆ 供应流程",
        "  宁德时代生产计划 → 法恩莱特接收订单 → 自动配料生产 →",
        "  质量检验 → 管道输送/槽车运输 → 宁德时代产线使用",
        "",
        "✓ 响应时间指标",
        "  订单接收：<5 分钟（系统自动接收）",
        "  生产启动：<30 分钟（自动化产线）",
        "  质量检验：<30 分钟（在线检测）",
        "  配送到达：<60 分钟（管道输送/槽车）",
        "  总响应时间：<2 小时（行业领先）",
        "",
        "◆ 库存策略",
        "  原材料库存：7-15 天安全库存",
        "  成品库存：零库存（按订单生产）",
        "  宁德时代库存：代管模式，用后结算"
    ])
    
    # ========== 第四部分：投资预算与回报 ==========
    print("添加第四部分...")
    add_section_slide(prs, "第四部分\n投资预算与回报")
    
    # 投资预算
    add_table_slide(prs, "单座撬装工厂投资预算", [
        ["类别", "金额（万元）", "占比"],
        ["固定资产投资", "4200-6050", "60%"],
        ["撬装设备", "2500-3500", "-"],
        ["管道输送系统", "300-500", "-"],
        ["土建工程", "500-800", "-"],
        ["自控系统", "300-400", "-"],
        ["无形资产投资", "180-310", "3%"],
        ["流动资金", "1500-2300", "27%"],
        ["预备费", "500-700", "10%"],
        ["总投资", "6380-9360", "100%"],
        [""],
        ["建议取值", "8000 万元/座（含 10% 预备费）", ""]
    ])
    
    # 投资回报测算
    add_table_slide(prs, "投资回报测算", [
        ["项目", "数值", "备注"],
        ["单座工厂产能", "5-10GWh", "配套能力"],
        ["单 GWh 电解液价值", "8000-10000 万元", "市场价格"],
        ["单座工厂年收入", "4-10 亿元", "满产状态"],
        ["产能利用率", "80%", "保守估计"],
        ["实际年收入", "3.2-8 亿元", "-"],
        ["年成本", "约 4 亿元", "原材料 70-80%"],
        ["年毛利", "1 亿元", "毛利率 20%"],
        ["年净利", "7000 万元", "净利率 14%"],
        ["投资回收期", "2-3 年", "保守估计"],
        [""],
        ["多工厂规划", "2027 年 2 座（1.6 亿），2029-30 年 4-6 座（3.2-4.8 亿）", ""]
    ])
    
    # ========== 第五部分：实施路径 ==========
    print("添加第五部分...")
    add_section_slide(prs, "第五部分\n实施路径")
    
    # 实施路径
    add_content_slide(prs, "实施路径 - 四个阶段", [
        "◆ 阶段一：验证期（2026 年 Q1-Q2）",
        "  完成撬装工厂设计方案（2 周）",
        "  与宁德时代对接供应商认证（4 周）",
        "  选址山东济宁或河南洛阳（2 周）",
        "  启动示范工厂建设（8 周）",
        "  里程碑：M6（第 22 周）试生产成功",
        "",
        "✓ 阶段二：落地期（2026 年 Q3-Q4）",
        "  通过宁德时代验厂（4 周）",
        "  小批量供货 0.5-1GWh（8 周）",
        "  验证 JIT 供应流程",
        "  里程碑：M10（第 20 周）年度供货目标达成",
        "",
        "◆ 阶段三：扩张期（2027 年）",
        "  启动第二座工厂建设（Q1）",
        "  第二座工厂投产（Q3）",
        "  产能提升至 10GWh（Q4）",
        "",
        "✓ 阶段四：全国布局期（2028-2030 年）",
        "  覆盖宁德时代全国基地",
        "  产能达到 50GWh",
        "  拓展海外市场"
    ])
    
    # 风险评估
    add_table_slide(prs, "风险评估与应对策略", [
        ["风险类别", "发生概率", "影响程度", "应对策略"],
        ["政策风险", "中", "高", "设备租赁模式、合资公司"],
        ["客户风险", "低", "高", "长期协议、拓展其他客户"],
        ["技术风险", "中", "中", "先槽车运输，逐步验证管道"],
        ["资金风险", "中", "中", "引入战投、绿色金融"],
        ["安全风险", "低", "高", "本质安全设计、HSE 体系"],
        ["环保风险", "低", "高", "清洁工艺、完善环保设施"],
        ["竞争风险", "高", "中", "快速占领、技术壁垒"]
    ])
    
    # ========== 第六部分：核心竞争优势 ==========
    print("添加第六部分...")
    add_section_slide(prs, "第六部分\n核心竞争优势")
    
    # 核心竞争优势
    add_content_slide(prs, "核心竞争优势", [
        "◆ 优势一：先发优势",
        "  国内首个电解液撬装工厂",
        "  建立行业标准，提高进入门槛",
        "",
        "✓ 优势二：成本优势",
        "  投资成本降低 50%（0.8 亿 vs 1.5-2 亿）",
        "  运营成本降低 30%",
        "  建设周期缩短 60%（3-6 个月 vs 2-3 年）",
        "",
        "◆ 优势三：响应优势",
        "  总响应时间<2 小时",
        "  领先传统供应商 12-24 倍（24-48 小时）",
        "",
        "✓ 优势四：绿色优势",
        "  零碳工厂，配套光伏 + 储能",
        "  符合宁德时代零碳供应链要求",
        "",
        "◆ 优势五：灵活优势",
        "  可快速复制（3-6 个月投产）",
        "  可灵活迁移（撬装设备可移动）"
    ])
    
    # ========== 第七部分：行动计划 ==========
    print("添加第七部分...")
    add_section_slide(prs, "第七部分\n行动计划")
    
    # 下一步行动
    add_content_slide(prs, "下一步行动计划", [
        "◆ 立即启动（本周）",
        "  成立项目筹备组（3-5 人）",
        "  与宁德时代供应链部门建立联系",
        "  调研山东济宁、河南洛阳基地周边土地政策",
        "  对接撬装设备供应商",
        "",
        "✓ 2 周内完成",
        "  编制详细商业计划书",
        "  完成投资可行性分析",
        "  确定首个示范工厂选址（优先洛阳）",
        "",
        "◆ 1 个月内完成",
        "  启动供应商认证流程",
        "  签订土地租赁意向",
        "  启动设备采购招标",
        "  组建项目团队（10-15 人）"
    ])
    
    # 3-6 个月计划
    add_content_slide(prs, "3-6 个月行动计划", [
        "◆ 3 个月内完成",
        "  完成详细设计",
        "  完成设备采购招标",
        "  完成土地租赁协议",
        "  完成安评、环评审批",
        "  启动工厂建设",
        "",
        "✓ 6 个月内完成",
        "  完成工厂建设",
        "  完成设备安装",
        "  完成系统调试",
        "  试生产",
        "  通过宁德时代验厂",
        "",
        "◆ 交付物",
        "  竣工验收报告、设备调试报告",
        "  试生产报告、验厂通过证书"
    ])
    
    # 战略建议
    add_content_slide(prs, "战略建议", [
        "◆ 建议一：立即启动",
        "  时不我待，竞争对手可能也在布局",
        "  建议本周内成立项目筹备组",
        "  快速推进，抢占先机",
        "",
        "✓ 建议二：优先洛阳",
        "  洛阳基地进展更快（一期已满产）",
        "  河南政策支持力度大",
        "  建议首选洛阳建设示范工厂",
        "",
        "◆ 建议三：轻资产运营",
        "  采用设备租赁模式",
        "  土地租赁而非购置",
        "  降低初始投资，提高资金效率",
        "",
        "✓ 建议四：技术领先",
        "  采用最先进自动化技术",
        "  实现 L4 级无人化操作",
        "  建立技术壁垒",
        "",
        "◆ 建议五：合作共赢",
        "  与宁德时代深度绑定",
        "  考虑成立合资公司",
        "  利益共享，风险共担"
    ])
    
    # 预期成果
    add_table_slide(prs, "预期成果", [
        ["阶段", "时间", "工厂数量", "产能", "销售收入"],
        ["短期", "2026 年", "1 座", "5GWh", "5 亿元"],
        ["中期", "2027-28 年", "2 座", "10GWh", "10 亿元"],
        ["长期", "2029-30 年", "4-6 座", "20-30GWh", "20-30 亿元"],
        [""],
        ["资本路径", "", "", "", ""],
        ["2028 年", "Pre-IPO 轮融资", "", "", ""],
        ["2029-30 年", "启动 IPO 或并购重组", "", "", ""]
    ])
    
    # 关键成功因素
    add_content_slide(prs, "关键成功因素", [
        "◆ 技术因素",
        "  撬装设备模块化设计",
        "  自动化控制系统稳定性",
        "  产品质量一致性",
        "",
        "✓ 市场因素",
        "  宁德时代供应商认证进度",
        "  与宁德时代采购部门关系",
        "  签订长期供货协议（3-5 年）",
        "",
        "◆ 政策因素",
        "  危化品许可政策突破",
        "  园区准入政策支持",
        "  环保、安全生产政策",
        "",
        "✓ 资金因素",
        "  融资渠道畅通",
        "  资金成本可控",
        "  回款周期管理"
    ])
    
    # ========== 结论与结束页 ==========
    print("添加结论与结束页...")
    
    # 结论
    add_content_slide(prs, "结论", [
        "◆ 核心结论",
        "  ✅ 模式可行：宁德时代'厂中厂'模式已验证",
        "  ✅ 市场巨大：山东、河南基地需求 20-24 万吨/年",
        "  ✅ 时机成熟：两大基地 2025-2026 年陆续投产",
        "  ✅ 投资可控：单座工厂 8000 万元，回收期 2-3 年",
        "  ✅ 优势明显：先发、成本、响应、绿色、灵活五大优势",
        "",
        "✓ 最终建议",
        "  🚀 立即启动项目，本周成立筹备组",
        "  🎯 优先布局洛阳，建设示范工厂",
        "  💡 采用轻资产运营，快速复制扩张",
        "  🏆 打造行业标杆，成为宁德时代核心供应商"
    ])
    
    # 结束页
    print("添加结束页...")
    add_title_slide(
        prs,
        "谢谢！",
        "法恩莱特撬装工厂战略规划\nQ&A"
    )
    
    # 保存 PPT
    output_file = "法恩莱特撬装工厂战略规划_专业版.pptx"
    print(f"保存 PPT 到：{output_file}")
    prs.save(output_file)
    
    print(f"\n[OK] 专业 PPT 生成完成！")
    print(f"文件位置：{output_file}")
    print(f"总页数：{len(prs.slides)} 页")
    
    return output_file

if __name__ == "__main__":
    generate_ppt()

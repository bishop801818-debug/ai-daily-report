"""
法恩莱特撬装工厂战略规划 PPT 生成脚本
使用 python-pptx 库自动生成 PowerPoint 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, prs):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 设置副标题
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def add_content_slide(prs, title, content_items, image_path=None):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 添加内容
    if content_items:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_after = Pt(10)
            p.level = 0
            
            # 如果是子要点
            if item.startswith('  -'):
                p.level = 1
                p.text = item.replace('  -', '')
    
    return slide

def add_data_slide(prs, title, table_data, chart_title=None):
    """添加数据表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白版式
    
    # 设置标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 添加表格
    if table_data:
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows > 0 and cols > 0:
            table_shape = slide.shapes.add_table(
                rows, cols,
                Inches(0.5), Inches(1.2),
                Inches(9), Inches(0.8 * rows)
            )
            
            table = table_shape.table
            
            # 填充数据
            for i, row in enumerate(table_data):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = str(cell)
                    table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(14)
                    
                    # 设置表头样式
                    if i == 0:
                        table.cell(i, j).fill.solid()
                        table.cell(i, j).fill.fore_color.rgb = RGBColor(0, 51, 102)
                        table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        table.cell(i, j).text_frame.paragraphs[0].font.bold = True
    
    return slide

def add_comparison_slide(prs, title, comparison_data):
    """添加对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 设置标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 添加对比内容
    if comparison_data:
        y_position = 1.2
        for item in comparison_data:
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(9), Inches(0.6))
            textbox.text = item
            textbox.text_frame.paragraphs[0].font.size = Pt(18)
            y_position += 0.6
    
    return slide

def generate_ppt():
    """生成 PPT"""
    print("开始生成 PPT...")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置页面大小为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. 标题页
    print("添加标题页...")
    add_title_slide(
        prs,
        "法恩莱特撬装工厂战略规划",
        "宁德时代山东、河南基地就近配套项目\n2026 年 4 月"
    )
    
    # 2. 目录
    print("添加目录页...")
    add_content_slide(prs, "目录", [
        "战略机遇验证",
        "宁德时代贴厂模式调研",
        "山东、河南基地布局",
        "战略定位与目标",
        "技术方案",
        "投资预算与回报",
        "实施路径",
        "核心竞争优势",
        "下一步行动计划"
    ])
    
    # 3. 执行摘要
    print("添加执行摘要...")
    add_content_slide(prs, "执行摘要 - 核心结论", [
        "✅ 宁德时代贴厂模式已验证",
        "  - 重庆'厂中厂'模式 2025 年 6 月投产",
        "  - 电池下线 1 小时内装车",
        "  - 供应链响应从数周压缩到分钟级",
        "✅ 山东、河南基地真实存在且规模巨大",
        "  - 山东济宁：2026 年产能超 100GWh",
        "  - 河南洛阳：2026 年产能超 100GWh",
        "✅ 电解液需求巨大",
        "  - 两大基地年需求 20-24 万吨",
        "  - 市场空间 60-96 亿元/年",
        "✅ 撬装工厂模式完全可行",
        "  - 投资回收期 2-3 年",
        "  - 建议立即启动"
    ])
    
    # 4. 宁德时代贴厂模式先例
    print("添加贴厂模式案例...")
    add_content_slide(prs, "宁德时代贴厂模式 - 重庆'厂中厂'", [
        "项目概况",
        "  - 时间：2025 年 6 月 30 日正式投产",
        "  - 地点：重庆赛力斯超级工厂内部",
        "  - 产能：两条 CTP2.0 高端电池包产线",
        "核心数据",
        "  - 电池下线后不到 1 小时就能装车",
        "  - 从下达指令到产线调整只需不到 20 分钟",
        "  - 每天数千个零部件当天周转完毕",
        "  - 物流响应周期从数天/数周压缩到分钟级",
        "模式特点",
        "  - 物理空间融合：供应商车间直接嵌入整车厂内部",
        "  - 生产系统协同：共享同一套生产节奏和数据系统",
        "  - 零距离供应：取消中间物流环节",
        "  - 零时差响应：实时响应生产计划变化"
    ])
    
    # 5. 山东基地布局
    print("添加山东基地信息...")
    add_data_slide(prs, "宁德时代山东基地布局", [
        ["项目", "详情"],
        ["位置", "山东省济宁市兖州区"],
        ["投资方", "宁德时代全资子公司 - 山东时代"],
        ["投产时间", "2025 年 5 月 17 日一期投产"],
        ["战略定位", "北方首个基地，北方产能规划最大基地"],
        ["一期产能", "60GWh，占地 800 余亩，投资 140 亿元"],
        ["二期产能", "30GWh，2025-2026 年"],
        ["三期产能", "40GWh+，2026 年"],
        ["总产能 (2026)", "超 100GWh，总占地约 2000 亩"],
        ["总投资", "超 210 亿元"],
        ["建设标准", "灯塔工厂 + 零碳工厂"]
    ])
    
    # 6. 河南基地布局
    print("添加河南基地信息...")
    add_data_slide(prs, "宁德时代河南基地布局", [
        ["项目", "详情"],
        ["位置", "河南省洛阳市伊滨区"],
        ["开工时间", "2022 年 9 月一期开工"],
        ["一期投产", "2024 年 7 月电池包工厂，2025 年 4 月全线投产"],
        ["二期投产", "2025 年 9 月/10 月投产"],
        ["一期产能", "30GWh，占地 850 亩"],
        ["二期产能", "30GWh"],
        ["三四期", "施工中，2026 年投产"],
        ["总产能", "超 100GWh，总占地 3400 亩"],
        ["总投资", "超 300 亿元"],
        ["产业链配套", "26 个上下游项目签约落地"],
        ["年产值", "全部建成后超千亿元"]
    ])
    
    # 7. 电解液需求测算
    print("添加需求测算...")
    add_data_slide(prs, "电解液需求测算", [
        ["基地", "2026 年产能", "电解液年需求", "市场空间"],
        ["山东济宁", "100GWh", "10-12 万吨", "30-48 亿元/年"],
        ["河南洛阳", "100GWh", "10-12 万吨", "30-48 亿元/年"],
        ["合计", "200GWh", "20-24 万吨", "60-96 亿元/年"],
        [""],
        ["计算公式：电解液需求（吨）= 电池产能（GWh）× 1000-1200 吨/GWh"]
    ])
    
    # 8. 战略定位
    print("添加战略定位...")
    add_content_slide(prs, "战略定位", [
        "总体定位",
        "成为宁德时代山东、河南基地的核心电解液战略供应商，",
        "采用撬装工厂模式实现就近配套、JIT 供应，",
        "打造零距离、零时差、零碳排的供应链标杆",
        "",
        "核心价值主张 - 四个'零'",
        "  - 零距离：撬装工厂部署在宁德时代基地周边 50 公里范围内",
        "  - 零时差：管道输送 + 自动化产线，响应时间<2 小时",
        "  - 零碳排：配套光伏 + 储能，使用绿电生产",
        "  - 零库存：按订单生产，JIT 供应",
        "",
        "三大优势",
        "  - 轻资产：模块化设计，可快速复制、灵活迁移",
        "  - 快响应：建设周期 3-6 个月，远快于传统工厂 2-3 年",
        "  - 低成本：投资成本降低 50%，运营成本降低 30%"
    ])
    
    # 9. 战略目标
    print("添加战略目标...")
    add_content_slide(prs, "战略目标体系", [
        "短期目标（2026 年）",
        "  - 完成 1 座撬装工厂示范项目建设",
        "  - 实现 5GWh 配套产能（约 5000 吨/年）",
        "  - 通过宁德时代供应商认证",
        "  - 示范工厂投资 5000-8000 万元",
        "",
        "中期目标（2027-2028 年）",
        "  - 在山东济宁、河南洛阳各部署 1 座撬装工厂",
        "  - 总产能达到 20GWh 配套（约 2 万吨/年）",
        "  - 进入宁德时代核心供应商体系",
        "  - 年销售收入达到 15-20 亿元",
        "",
        "长期目标（2029-2030 年）",
        "  - 部署 4-6 座撬装工厂，覆盖宁德时代全国基地",
        "  - 总产能达到 50GWh 配套（约 5 万吨/年）",
        "  - 成为宁德时代前三大电解液供应商",
        "  - 年销售收入达到 40-50 亿元"
    ])
    
    # 10. 撬装工厂设计参数
    print("添加技术参数...")
    add_data_slide(prs, "撬装工厂设计参数", [
        ["参数项", "指标", "传统工厂对比"],
        ["占地面积", "800-1000㎡", "传统工厂 1/10"],
        ["建筑面积", "1000-1500㎡", "传统工厂 1/8"],
        ["建设周期", "3-6 个月", "传统工厂 1/3"],
        ["投资规模", "5000-8000 万元", "传统工厂 1/2"],
        ["产能规模", "5-10GWh 配套", "可灵活扩展"],
        ["人员配置", "10-15 人", "传统工厂 1/5"],
        ["自动化程度", "L4 级（无人化）", "行业领先"]
    ])
    
    # 11. 核心功能模块
    print("添加功能模块...")
    add_content_slide(prs, "撬装工厂核心功能模块", [
        "模块一：原料储存模块",
        "  - NMP 储罐、锂盐储罐、添加剂储罐",
        "  - 氮气保护系统、温湿度控制",
        "",
        "模块二：配料混合模块",
        "  - 高精度计量系统（精度±0.1%）",
        "  - 自动配料系统、真空搅拌系统",
        "",
        "模块三：过滤灌装模块",
        "  - 多级过滤系统（精度 1μm）",
        "  - 自动灌装系统、密封检测",
        "",
        "模块四：成品仓储模块",
        "  - 恒温恒湿仓库、自动立体仓库",
        "  - AGV 搬运系统、WMS 管理系统",
        "",
        "模块五：管道输送模块（可选）",
        "  - 不锈钢管道系统、计量泵组",
        "  - 直接输送到宁德时代产线"
    ])
    
    # 12. JIT 供应流程
    print("添加供应链流程...")
    add_content_slide(prs, "JIT 供应流程与响应时间", [
        "供应流程",
        "宁德时代生产计划 → 法恩莱特接收订单 → 自动配料生产 →",
        "质量检验 → 管道输送/槽车运输 → 宁德时代产线使用",
        "",
        "响应时间指标",
        "  - 订单接收：<5 分钟（系统自动接收）",
        "  - 生产启动：<30 分钟（自动化产线）",
        "  - 质量检验：<30 分钟（在线检测）",
        "  - 配送到达：<60 分钟（管道输送/槽车）",
        "  - 总响应时间：<2 小时（行业领先）",
        "",
        "库存策略",
        "  - 原材料库存：7-15 天安全库存",
        "  - 成品库存：零库存（按订单生产）",
        "  - 宁德时代库存：代管模式，用后结算"
    ])
    
    # 13. 投资预算
    print("添加投资预算...")
    add_data_slide(prs, "单座撬装工厂投资预算", [
        ["类别", "金额（万元）", "占比"],
        ["固定资产投资", "4200-6050", "60%"],
        ["  - 撬装设备", "2500-3500", "-"],
        ["  - 管道输送系统", "300-500", "-"],
        ["  - 土建工程", "500-800", "-"],
        ["  - 自控系统", "300-400", "-"],
        ["无形资产投资", "180-310", "3%"],
        ["  - 设计、安评、环评等", "180-310", "-"],
        ["流动资金", "1500-2300", "27%"],
        ["预备费", "500-700", "10%"],
        ["总投资", "6380-9360", "100%"],
        [""],
        ["建议取值：8000 万元/座（含 10% 预备费）"]
    ])
    
    # 14. 投资回报
    print("添加投资回报...")
    add_data_slide(prs, "投资回报测算", [
        ["项目", "数值", "备注"],
        ["单座工厂产能", "5-10GWh", "配套能力"],
        ["单 GWh 电解液价值", "8000-10000 万元", "市场价格"],
        ["单座工厂年收入", "4-10 亿元", "满产状态"],
        ["产能利用率", "80%", "保守估计"],
        ["实际年收入", "3.2-8 亿元", "-"],
        ["年成本", "约 4 亿元", "原材料 70-80%"],
        ["年毛利", "1 亿元", "毛利率 20%"],
        ["年净利", "7000 万元", "净利率 14%"],
        ["投资回收期", "2-3 年", "保守估计"],
        [""],
        ["多工厂规划：2027 年 2 座（1.6 亿），2029-30 年 4-6 座（3.2-4.8 亿）"]
    ])
    
    # 15. 实施路径
    print("添加实施路径...")
    add_content_slide(prs, "实施路径 - 四个阶段", [
        "阶段一：验证期（2026 年 Q1-Q2）",
        "  - 完成撬装工厂设计方案（2 周）",
        "  - 与宁德时代对接供应商认证（4 周）",
        "  - 选址山东济宁或河南洛阳（2 周）",
        "  - 启动示范工厂建设（8 周）",
        "  - 里程碑：M6（第 22 周）试生产成功",
        "",
        "阶段二：落地期（2026 年 Q3-Q4）",
        "  - 通过宁德时代验厂（4 周）",
        "  - 小批量供货 0.5-1GWh（8 周）",
        "  - 验证 JIT 供应流程",
        "  - 里程碑：M10（第 20 周）年度供货目标达成",
        "",
        "阶段三：扩张期（2027 年）",
        "  - 启动第二座工厂建设（Q1）",
        "  - 第二座工厂投产（Q3）",
        "  - 产能提升至 10GWh（Q4）",
        "",
        "阶段四：全国布局期（2028-2030 年）",
        "  - 覆盖宁德时代全国基地",
        "  - 产能达到 50GWh",
        "  - 拓展海外市场"
    ])
    
    # 16. 风险评估
    print("添加风险评估...")
    add_data_slide(prs, "风险评估与应对策略", [
        ["风险类别", "发生概率", "影响程度", "应对策略"],
        ["政策风险", "中", "高", "设备租赁模式、合资公司"],
        ["客户风险", "低", "高", "长期协议、拓展其他客户"],
        ["技术风险", "中", "中", "先槽车运输，逐步验证管道"],
        ["资金风险", "中", "中", "引入战投、绿色金融"],
        ["安全风险", "低", "高", "本质安全设计、HSE 体系"],
        ["环保风险", "低", "高", "清洁工艺、完善环保设施"],
        ["竞争风险", "高", "中", "快速占领、技术壁垒"]
    ])
    
    # 17. 核心竞争优势
    print("添加竞争优势...")
    add_content_slide(prs, "核心竞争优势", [
        "优势一：先发优势",
        "  - 国内首个电解液撬装工厂",
        "  - 建立行业标准，提高进入门槛",
        "",
        "优势二：成本优势",
        "  - 投资成本降低 50%（0.8 亿 vs 1.5-2 亿）",
        "  - 运营成本降低 30%",
        "  - 建设周期缩短 60%（3-6 个月 vs 2-3 年）",
        "",
        "优势三：响应优势",
        "  - 总响应时间<2 小时",
        "  - 领先传统供应商 12-24 倍（24-48 小时）",
        "",
        "优势四：绿色优势",
        "  - 零碳工厂，配套光伏 + 储能",
        "  - 符合宁德时代零碳供应链要求",
        "",
        "优势五：灵活优势",
        "  - 可快速复制（3-6 个月投产）",
        "  - 可灵活迁移（撬装设备可移动）"
    ])
    
    # 18. 下一步行动
    print("添加行动计划...")
    add_content_slide(prs, "下一步行动计划", [
        "立即启动（本周）",
        "  - ✅ 成立项目筹备组（3-5 人）",
        "  - ✅ 与宁德时代供应链部门建立联系",
        "  - ✅ 调研山东济宁、河南洛阳基地周边土地政策",
        "  - ✅ 对接撬装设备供应商",
        "",
        "2 周内完成",
        "  - 编制详细商业计划书",
        "  - 完成投资可行性分析",
        "  - 确定首个示范工厂选址（优先洛阳）",
        "",
        "1 个月内完成",
        "  - 启动供应商认证流程",
        "  - 签订土地租赁意向",
        "  - 启动设备采购招标",
        "  - 组建项目团队（10-15 人）"
    ])
    
    # 19. 3 个月行动计划
    print("添加 3-6 个月计划...")
    add_content_slide(prs, "3-6 个月行动计划", [
        "3 个月内完成",
        "  - 完成详细设计",
        "  - 完成设备采购招标",
        "  - 完成土地租赁协议",
        "  - 完成安评、环评审批",
        "  - 启动工厂建设",
        "",
        "6 个月内完成",
        "  - 完成工厂建设",
        "  - 完成设备安装",
        "  - 完成系统调试",
        "  - 试生产",
        "  - 通过宁德时代验厂",
        "",
        "交付物",
        "  - 竣工验收报告、设备调试报告",
        "  - 试生产报告、验厂通过证书"
    ])
    
    # 20. 战略建议
    print("添加战略建议...")
    add_content_slide(prs, "战略建议", [
        "建议一：立即启动",
        "  - 时不我待，竞争对手可能也在布局",
        "  - 建议本周内成立项目筹备组",
        "  - 快速推进，抢占先机",
        "",
        "建议二：优先洛阳",
        "  - 洛阳基地进展更快（一期已满产）",
        "  - 河南政策支持力度大",
        "  - 建议首选洛阳建设示范工厂",
        "",
        "建议三：轻资产运营",
        "  - 采用设备租赁模式",
        "  - 土地租赁而非购置",
        "  - 降低初始投资，提高资金效率",
        "",
        "建议四：技术领先",
        "  - 采用最先进自动化技术",
        "  - 实现 L4 级无人化操作",
        "  - 建立技术壁垒",
        "",
        "建议五：合作共赢",
        "  - 与宁德时代深度绑定",
        "  - 考虑成立合资公司",
        "  - 利益共享，风险共担"
    ])
    
    # 21. 预期成果
    print("添加预期成果...")
    add_data_slide(prs, "预期成果", [
        ["阶段", "时间", "工厂数量", "产能", "销售收入"],
        ["短期", "2026 年", "1 座", "5GWh", "5 亿元"],
        ["中期", "2027-28 年", "2 座", "10GWh", "10 亿元"],
        ["长期", "2029-30 年", "4-6 座", "20-30GWh", "20-30 亿元"],
        [""],
        ["资本路径", "", "", "", ""],
        ["2028 年", "Pre-IPO 轮融资", "", "", ""],
        ["2029-30 年", "启动 IPO 或并购重组", "", "", ""]
    ])
    
    # 22. 关键成功因素
    print("添加成功因素...")
    add_content_slide(prs, "关键成功因素", [
        "技术因素",
        "  - 撬装设备模块化设计",
        "  - 自动化控制系统稳定性",
        "  - 产品质量一致性",
        "",
        "市场因素",
        "  - 宁德时代供应商认证进度",
        "  - 与宁德时代采购部门关系",
        "  - 签订长期供货协议（3-5 年）",
        "",
        "政策因素",
        "  - 危化品许可政策突破",
        "  - 园区准入政策支持",
        "  - 环保、安全生产政策",
        "",
        "资金因素",
        "  - 融资渠道畅通",
        "  - 资金成本可控",
        "  - 回款周期管理"
    ])
    
    # 23. 结论
    print("添加结论页...")
    add_content_slide(prs, "结论", [
        "核心结论",
        "  ✅ 模式可行：宁德时代'厂中厂'模式已验证",
        "  ✅ 市场巨大：山东、河南基地需求 20-24 万吨/年",
        "  ✅ 时机成熟：两大基地 2025-2026 年陆续投产",
        "  ✅ 投资可控：单座工厂 8000 万元，回收期 2-3 年",
        "  ✅ 优势明显：先发、成本、响应、绿色、灵活五大优势",
        "",
        "最终建议",
        "  🚀 立即启动项目，本周成立筹备组",
        "  🎯 优先布局洛阳，建设示范工厂",
        "  💡 采用轻资产运营，快速复制扩张",
        "  🏆 打造行业标杆，成为宁德时代核心供应商"
    ])
    
    # 24. 结束页
    print("添加结束页...")
    add_title_slide(
        prs,
        "谢谢！",
        "法恩莱特撬装工厂战略规划\nQ&A"
    )
    
    # 保存 PPT
    output_file = "法恩莱特撬装工厂战略规划.pptx"
    print(f"保存 PPT 到：{output_file}")
    prs.save(output_file)
    
    print(f"\n[OK] PPT 生成完成！")
    print(f"文件位置：{output_file}")
    print(f"总页数：{len(prs.slides)} 页")
    
    return output_file

if __name__ == "__main__":
    generate_ppt()
